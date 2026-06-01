"""外部 .pptx 模板加载器

支持加载外部 .pptx 模板文件，自动推断布局映射，提取主题颜色/字体。
"""

from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn

from ppt_maker.models import SlideType, ThemeModel, ColorScheme, FontScheme


class ExternalTemplate:
    """外部 .pptx 模板处理器"""

    def __init__(self, path: str, layout_map: dict | None = None):
        """
        Args:
            path: .pptx 模板文件路径
            layout_map: 手动布局映射 {SlideType.name -> layout_index}
        """
        self.prs = Presentation(path)
        self._manual_map = layout_map or {}
        self._auto_map: dict[str, int] = {}
        self._analyze_layouts()

    def _analyze_layouts(self):
        """分析模板中的布局，按名称/占位符推断类型"""
        for i, layout in enumerate(self.prs.slide_layouts):
            name = layout.name.lower()

            # 先尝试按名称推断
            if "title slide" in name:
                self._auto_map[SlideType.TITLE.name] = i
            elif "section" in name or "section header" in name:
                self._auto_map[SlideType.SECTION.name] = i
            elif "two content" in name or "two-column" in name or "2 content" in name:
                self._auto_map[SlideType.TWO_COLUMN.name] = i
            elif "picture with caption" in name:
                self._auto_map[SlideType.IMAGE_TEXT.name] = i
            elif "content with caption" in name:
                self._auto_map[SlideType.CONTENT.name] = i
            elif "blank" in name:
                self._auto_map[SlideType.BLANK.name] = i
            elif "title and content" in name or "title and text" in name:
                self._auto_map[SlideType.CONTENT.name] = i
            else:
                # 按占位符结构推断
                self._infer_by_placeholders(layout, i)

        # 确保 CONTENT 有默认值
        if SlideType.CONTENT.name not in self._auto_map:
            # 使用 title and content 布局或第一个有 title 占位符的
            for i, layout in enumerate(self.prs.slide_layouts):
                if any(ph.placeholder_format.idx == 1 for ph in layout.placeholders):
                    self._auto_map[SlideType.CONTENT.name] = i
                    break

        # 确保 TITLE 有默认值
        if SlideType.TITLE.name not in self._auto_map:
            self._auto_map[SlideType.TITLE.name] = 0

        # 确保 BLANK 有默认值
        if SlideType.BLANK.name not in self._auto_map:
            # 找最后面的布局通常是 blank
            self._auto_map[SlideType.BLANK.name] = len(self.prs.slide_layouts) - 1

    def _infer_by_placeholders(self, layout, index: int):
        """按占位符结构推断布局类型"""
        ph_indices = [ph.placeholder_format.idx for ph in layout.placeholders]

        if 1 in ph_indices and 2 in ph_indices and 3 not in ph_indices:
            # title + subtitle -> TITLE
            if SlideType.TITLE.name not in self._auto_map:
                self._auto_map[SlideType.TITLE.name] = index
        elif 1 in ph_indices and 2 in ph_indices:
            # title + body -> CONTENT
            if SlideType.CONTENT.name not in self._auto_map:
                self._auto_map[SlideType.CONTENT.name] = index

    def get_layout_for_type(self, slide_type: SlideType):
        """根据 SlideType 获取对应的 python-pptx SlideLayout

        优先使用手动映射，其次自动推断。
        """
        # 手动映射优先
        if slide_type.value in self._manual_map:
            idx = self._manual_map[slide_type.value]
            return self.prs.slide_layouts[idx]

        # 自动映射
        if slide_type.name in self._auto_map:
            idx = self._auto_map[slide_type.name]
            return self.prs.slide_layouts[idx]

        # 回退: CONTENT -> title and content 布局, 其他 -> blank
        if slide_type in (SlideType.CONTENT, SlideType.TITLE, SlideType.SECTION):
            # 尝试找带 title 的布局
            for layout in self.prs.slide_layouts:
                if any(ph.placeholder_format.idx == 1 for ph in layout.placeholders):
                    return layout

        # 最终回退到第一个布局
        return self.prs.slide_layouts[0]

    def extract_theme(self) -> ThemeModel:
        """从模板中提取主题颜色和字体

        读取 slide_master 中的主题 XML。
        """
        try:
            color_scheme = self._extract_colors()
        except Exception:
            color_scheme = ColorScheme()

        try:
            font_scheme = self._extract_fonts()
        except Exception:
            font_scheme = FontScheme()

        return ThemeModel(
            name="external_template",
            display_name="外部模板",
            color_scheme=color_scheme,
            font_scheme=font_scheme,
        )

    def _extract_colors(self) -> ColorScheme:
        """从模板 XML 提取颜色方案"""
        master = self.prs.slide_masters[0]
        # 尝试从主题 XML 中提取
        try:
            theme_element = master.element.find(
                ".//" + qn("a:clrScheme")
            )
            if theme_element is None:
                # 尝试从 slide_layout 的主题获取
                return ColorScheme()

            colors = {}
            for child in theme_element:
                tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                srgb = child.find(qn("a:srgbClr"))
                if srgb is not None:
                    colors[tag] = f"#{srgb.get('val', '000000')}"

            return ColorScheme(
                primary=colors.get("dk1", "#000000"),
                secondary=colors.get("lt1", "#FFFFFF"),
                accent=colors.get("accent1", "#4472C4"),
                background=colors.get("lt2", "#F5F5F5"),
                text_dark=colors.get("dk1", "#000000"),
                text_light=colors.get("lt1", "#FFFFFF"),
                text_secondary=colors.get("dk2", "#44546A"),
                divider=colors.get("lt2", "#E7E6E6"),
            )
        except Exception:
            return ColorScheme()

    def _extract_fonts(self) -> FontScheme:
        """从模板 XML 提取字体方案"""
        try:
            master = self.prs.slide_masters[0]
            font_element = master.element.find(
                ".//" + qn("a:fontScheme")
            )
            if font_element is None:
                return FontScheme()

            major = font_element.find(qn("a:majorFont"))
            minor = font_element.find(qn("a:minorFont"))

            title_font = "微软雅黑"
            title_font_latin = "Calibri"
            body_font = "微软雅黑"
            body_font_latin = "Calibri"

            if major is not None:
                latin = major.find(qn("a:latin"))
                ea = major.find(qn("a:ea"))
                if latin is not None:
                    title_font_latin = latin.get("typeface", "Calibri")
                if ea is not None:
                    title_font = ea.get("typeface", "微软雅黑")

            if minor is not None:
                latin = minor.find(qn("a:latin"))
                ea = minor.find(qn("a:ea"))
                if latin is not None:
                    body_font_latin = latin.get("typeface", "Calibri")
                if ea is not None:
                    body_font = ea.get("typeface", "微软雅黑")

            return FontScheme(
                title_font=title_font,
                title_font_latin=title_font_latin,
                subtitle_font=title_font,
                subtitle_font_latin=title_font_latin,
                body_font=body_font,
                body_font_latin=body_font_latin,
                caption_font=body_font,
                caption_font_latin=body_font_latin,
            )
        except Exception:
            return FontScheme()

    @property
    def presentation(self) -> Presentation:
        """返回底层 python-pptx Presentation 对象"""
        return self.prs
