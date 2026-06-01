"""CLI 入口 - Typer 命令行界面"""

import os
import typer
from typing import Optional
from pathlib import Path
from rich import print as rprint
from rich.console import Console

from ppt_maker import create_from_outline, create_from_ai, list_builtin_themes
from ppt_maker.config import Config

console = Console()

app = typer.Typer(
    name="ppt-maker",
    help="专业 PPT 创建工具 - 支持大纲输入和 AI 辅助生成",
    rich_markup_mode="rich",
)


def _load_config(config_path: str | None) -> Config:
    """加载配置文件"""
    if config_path:
        return Config.from_file(config_path)
    return Config()


@app.command()
def create(
    input: str = typer.Argument(..., help="输入文件路径 (Markdown/JSON/YAML)"),
    output: str = typer.Option("output.pptx", "-o", "--output", help="输出文件路径"),
    theme: str = typer.Option("default", "-t", "--theme", help="主题名称"),
    template: Optional[str] = typer.Option(None, "--template", help="外部 .pptx 模板路径"),
    config: Optional[str] = typer.Option(None, "-c", "--config", help="配置文件路径"),
):
    """从大纲文件创建 PPT"""
    if not Path(input).exists():
        rprint(f"[red]错误: 文件不存在: {input}[/red]")
        raise typer.Exit(code=1)

    if template and not Path(template).exists():
        rprint(f"[red]错误: 模板文件不存在: {template}[/red]")
        raise typer.Exit(code=1)

    # 加载配置
    cfg = _load_config(config)
    if not template and cfg.template_path:
        template = cfg.template_path
    if theme == "default" and cfg.theme_name != "default":
        theme = cfg.theme_name

    rprint(f"正在解析大纲: {input}")
    try:
        output_path = create_from_outline(
            input_path=input,
            output_path=output,
            theme_name=theme,
            template_path=template,
        )
        rprint(f"[green]PPT 已生成: {output_path}[/green]")
    except Exception as e:
        rprint(f"[red]错误: {e}[/red]")
        raise typer.Exit(code=1)


@app.command()
def generate(
    topic: str = typer.Argument(..., help="演示文稿主题"),
    output: str = typer.Option("output.pptx", "-o", "--output", help="输出文件路径"),
    slides: int = typer.Option(10, "-n", "--slides", help="幻灯片数量"),
    theme: str = typer.Option("default", "-t", "--theme", help="主题名称"),
    language: str = typer.Option("zh-CN", "-l", "--language", help="语言"),
    model: Optional[str] = typer.Option(None, "--model", help="AI 模型"),
    base_url: Optional[str] = typer.Option(None, "--base-url", help="API 端点 URL (支持 deepseek/moonshot/ollama)"),
    instructions: Optional[str] = typer.Option(None, "-i", "--instructions", help="附加生成指令"),
    template: Optional[str] = typer.Option(None, "--template", help="外部 .pptx 模板路径"),
    config: Optional[str] = typer.Option(None, "-c", "--config", help="配置文件路径"),
):
    """使用 AI 生成 PPT"""
    # 加载配置
    cfg = _load_config(config)

    # 配置合并: CLI 参数 > 配置文件 > 默认值
    resolved_model = model or cfg.ai_model
    resolved_base_url = base_url or cfg.ai_base_url
    if not template and cfg.template_path:
        template = cfg.template_path
    if theme == "default" and cfg.theme_name != "default":
        theme = cfg.theme_name

    api_key = os.environ.get("OPENAI_API_KEY")
    if not api_key:
        rprint("[red]错误: 请设置 OPENAI_API_KEY 环境变量[/red]")
        raise typer.Exit(code=1)

    if template and not Path(template).exists():
        rprint(f"[red]错误: 模板文件不存在: {template}[/red]")
        raise typer.Exit(code=1)

    rprint(f"正在使用 AI 生成 PPT: [bold]{topic}[/bold]")
    rprint(f"模型: {resolved_model}, 页数: {slides}")

    try:
        output_path = create_from_ai(
            topic=topic,
            output_path=output,
            slide_count=slides,
            theme_name=theme,
            language=language,
            instructions=instructions or "",
            ai_model=resolved_model,
            api_key=api_key,
            base_url=resolved_base_url,
            template_path=template,
        )
        rprint(f"[green]PPT 已生成: {output_path}[/green]")
    except Exception as e:
        rprint(f"[red]错误: {e}[/red]")
        raise typer.Exit(code=1)


@app.command()
def themes():
    """列出所有可用主题"""
    from ppt_maker.template.builtin import load_builtin_theme

    available = list_builtin_themes()
    rprint("[bold]可用主题:[/bold]")
    for name in available:
        theme = load_builtin_theme(name)
        rprint(f"  [cyan]{name}[/cyan]: {theme.display_name}")


@app.command()
def inspect(
    file: str = typer.Argument(..., help=".pptx 文件路径"),
):
    """检查 .pptx 文件信息 (布局、幻灯片数量等)"""
    if not Path(file).exists():
        rprint(f"[red]错误: 文件不存在: {file}[/red]")
        raise typer.Exit(code=1)

    try:
        from pptx import Presentation
        from pptx.util import Emu
        from ppt_maker.template.external import ExternalTemplate

        prs = Presentation(file)
        rprint(f"[bold]文件:[/bold] {file}")
        rprint(f"[bold]幻灯片数量:[/bold] {len(prs.slides)}")

        # 转换 EMU 为 cm
        w_cm = round(prs.slide_width / 914400, 1)
        h_cm = round(prs.slide_height / 914400, 1)
        rprint(f"[bold]尺寸:[/bold] {w_cm} x {h_cm} cm")
        rprint(f"[bold]布局数量:[/bold] {len(prs.slide_layouts)}")

        rprint("\n[bold]布局列表:[/bold]")
        for i, layout in enumerate(prs.slide_layouts):
            ph_count = len(layout.placeholders)
            rprint(f"  [{i}] {layout.name} ({ph_count} 占位符)")

        # 幻灯片概览
        rprint(f"\n[bold]幻灯片概览:[/bold]")
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text[:30]
                    if text.strip():
                        texts.append(text.strip())
            preview = " | ".join(texts[:2]) if texts else "(无文本)"
            rprint(f"  页{i+1}: {preview}")

        # 提取主题信息
        ext = ExternalTemplate(file)
        theme = ext.extract_theme()
        rprint(f"\n[bold]主题颜色:[/bold]")
        rprint(f"  主色: [#{theme.color_scheme.primary.lstrip('#')}]{theme.color_scheme.primary}[/#{theme.color_scheme.primary.lstrip('#')}]")
        rprint(f"  背景色: {theme.color_scheme.background}")
        rprint(f"  正文色: {theme.color_scheme.text_dark}")

        rprint(f"\n[bold]主题字体:[/bold]")
        rprint(f"  标题: {theme.font_scheme.title_font} / {theme.font_scheme.title_font_latin}")
        rprint(f"  正文: {theme.font_scheme.body_font} / {theme.font_scheme.body_font_latin}")

    except Exception as e:
        rprint(f"[red]错误: 无法读取文件: {e}[/red]")
        raise typer.Exit(code=1)


if __name__ == "__main__":
    app()
