"""图表渲染器"""

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor

from ppt_maker.models import ChartElement, ChartType, LayoutZone, ThemeModel


# ChartType -> XL_CHART_TYPE 映射
CHART_TYPE_MAP = {
    ChartType.BAR: XL_CHART_TYPE.COLUMN_CLUSTERED,
    ChartType.BAR_STACKED: XL_CHART_TYPE.COLUMN_STACKED,
    ChartType.LINE: XL_CHART_TYPE.LINE_MARKERS,
    ChartType.PIE: XL_CHART_TYPE.PIE,
    ChartType.AREA: XL_CHART_TYPE.AREA,
    ChartType.SCATTER: XL_CHART_TYPE.XY_SCATTER,
}

# 主题色板 (用于系列着色)
SERIES_COLORS = [
    "4472C4",  # 蓝
    "ED7D31",  # 橙
    "A5A5A5",  # 灰
    "FFC000",  # 金
    "5B9BD5",  # 浅蓝
    "70AD47",  # 绿
    "264478",  # 深蓝
    "9B59B6",  # 紫
]


class ChartRenderer:
    """图表渲染器"""

    def __init__(self, theme: ThemeModel):
        self.theme = theme

    def render_chart(self, slide, element: ChartElement, zone: LayoutZone):
        """渲染图表元素到幻灯片"""
        chart_type_enum = CHART_TYPE_MAP.get(
            element.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED
        )

        chart_data = self._build_chart_data(element)

        left = Emu(int(zone.x))
        top = Emu(int(zone.y))
        width = Emu(int(zone.width))
        height = Emu(int(zone.height))

        graphic_frame = slide.shapes.add_chart(
            chart_type_enum, left, top, width, height, chart_data
        )
        chart = graphic_frame.chart

        self._style_chart(chart, element)

    def _build_chart_data(self, element: ChartElement) -> CategoryChartData:
        """构建图表数据"""
        chart_data = CategoryChartData()
        chart_data.categories = element.categories

        for series in element.series:
            chart_data.add_series(series.name, tuple(series.values))

        return chart_data

    def _style_chart(self, chart, element: ChartElement):
        """应用图表样式"""
        # 图表标题
        if element.title:
            chart.has_title = True
            chart.chart_title.text_frame.paragraphs[0].text = element.title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
            try:
                chart.chart_title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(
                    self.theme.color_scheme.text_dark.lstrip("#")
                )
            except Exception:
                pass
        else:
            chart.has_title = False

        # 图例
        chart.has_legend = element.show_legend
        if element.show_legend:
            chart.legend.include_in_layout = False
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            try:
                chart.legend.font.size = Pt(10)
            except Exception:
                pass

        # 数据标签
        if element.show_data_labels:
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.font.size = Pt(9)
            data_labels.number_format = '#,##0'
            if element.chart_type == ChartType.PIE:
                data_labels.show_percentage = True
                data_labels.show_category_name = True
                data_labels.show_value = False
            else:
                data_labels.show_value = True
                data_labels.show_category_name = False

        # 系列颜色
        self._apply_series_colors(chart, element)

    def _apply_series_colors(self, chart, element: ChartElement):
        """为图表系列应用主题色"""
        plot = chart.plots[0]

        for i, series in enumerate(plot.series):
            color_hex = SERIES_COLORS[i % len(SERIES_COLORS)]
            try:
                fill = series.format.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor.from_string(color_hex)
            except Exception:
                pass

            # 饼图每个点单独着色
            if element.chart_type == ChartType.PIE:
                for j, point in enumerate(series.points):
                    color_hex = SERIES_COLORS[j % len(SERIES_COLORS)]
                    try:
                        fill = point.format.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor.from_string(color_hex)
                    except Exception:
                        pass
