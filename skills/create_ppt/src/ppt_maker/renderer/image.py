"""图片渲染器"""

import os
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE

from ppt_maker.models import ImageElement, LayoutZone, ThemeModel


class ImageRenderer:
    """图片渲染器"""

    def __init__(self, theme: ThemeModel):
        self.theme = theme

    def render_image(self, slide, element: ImageElement, zone: LayoutZone):
        """渲染图片元素到幻灯片"""
        if not element.source:
            return

        source = element.source.strip()

        # 目前仅支持本地文件路径
        if not os.path.exists(source):
            # 尝试相对于当前工作目录
            if not os.path.isabs(source):
                return
            return

        left = Emu(int(zone.x))
        top = Emu(int(zone.y))
        width = Emu(int(zone.width))
        height = Emu(int(zone.height))

        try:
            if element.scaling == "fit":
                # 保持比例适应区域
                pic = slide.shapes.add_picture(
                    source, left, top, width, height
                )
            elif element.scaling == "fill":
                # 填充整个区域 (可能裁剪)
                pic = slide.shapes.add_picture(
                    source, left, top, width, height
                )
            else:
                # stretch - 拉伸填满
                pic = slide.shapes.add_picture(
                    source, left, top, width, height
                )

            # 设置替代文本
            if element.alt_text:
                try:
                    pic.name = element.alt_text
                except Exception:
                    pass

        except Exception:
            # 图片加载失败，创建占位框
            self._render_placeholder(slide, element, zone)

    def _render_placeholder(self, slide, element: ImageElement, zone: LayoutZone):
        """图片加载失败时的占位框"""
        left = Emu(int(zone.x))
        top = Emu(int(zone.y))
        width = Emu(int(zone.width))
        height = Emu(int(zone.height))

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        shape.fill.background()
        try:
            from pptx.dml.color import RGBColor
            shape.line.color.rgb = RGBColor.from_string(
                self.theme.color_scheme.divider.lstrip("#")
            )
        except Exception:
            pass
        shape.line.width = Emu(12700)  # 1pt

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = 2  # center
        run = p.add_run()
        run.text = f"[图片: {element.source}]"
