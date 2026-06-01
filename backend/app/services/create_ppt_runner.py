"""Runner helpers for the bundled create_ppt skill."""

from __future__ import annotations

import importlib.util
import json
import re
import sys
import uuid
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

from fastapi import HTTPException


_MODULE_NAME = "lobster_create_ppt_pipeline"


def _lobster_root() -> Path:
    return Path(__file__).resolve().parents[3]


def _pipeline_script() -> Path:
    return _lobster_root() / "skills" / "create_ppt" / "scripts" / "create_ppt_pipeline.py"


def default_create_ppt_output_dir() -> Path:
    return _lobster_root() / "_lobster_runtime" / "create_ppt" / "runs"


def safe_create_ppt_name(value: str, default: str = "presentation") -> str:
    text = (value or "").strip() or default
    text = re.sub(r'[<>:"/\\|?*\x00-\x1f]+', "_", text)
    text = re.sub(r"\s+", "_", text).strip("._ ")
    return (text or default)[:80]


def create_ppt_run_dir(topic: str) -> Path:
    output_dir = default_create_ppt_output_dir()
    output_dir.mkdir(parents=True, exist_ok=True)
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    base_name = safe_create_ppt_name(topic or "presentation")
    run_dir = output_dir / f"run_{stamp}_{base_name}"
    n = 1
    while run_dir.exists():
        n += 1
        run_dir = output_dir / f"run_{stamp}_{base_name}_{n:02d}"
    run_dir.mkdir(parents=True, exist_ok=True)
    return run_dir


def _load_pipeline_module():
    script = _pipeline_script()
    if not script.exists():
        raise HTTPException(status_code=503, detail=f"未找到 PPT 生成脚本: {script}")
    old = sys.modules.get(_MODULE_NAME)
    if old is not None:
        return old
    spec = importlib.util.spec_from_file_location(_MODULE_NAME, script)
    if not spec or not spec.loader:
        raise HTTPException(status_code=503, detail="无法加载 PPT 生成模块")
    mod = importlib.util.module_from_spec(spec)
    sys.modules[_MODULE_NAME] = mod
    spec.loader.exec_module(mod)
    return mod


def run_create_ppt_sync(payload: Dict[str, Any]) -> Dict[str, Any]:
    data = dict(payload or {})
    if not str(data.get("output_dir") or "").strip():
        data["output_dir"] = str(default_create_ppt_output_dir())
    mod = _load_pipeline_module()
    result = mod.run_pipeline(data)
    if not isinstance(result, dict) or not result.get("ok"):
        raise HTTPException(status_code=500, detail=result or "PPT 生成失败")
    return result


def create_fullpage_image_ppt(
    *,
    image_paths: List[str],
    output_path: str,
    notes: Optional[List[str]] = None,
) -> str:
    if not image_paths:
        raise HTTPException(status_code=400, detail="没有可用于生成 PPT 的图片")
    try:
        from pptx import Presentation
    except Exception as exc:
        raise HTTPException(status_code=503, detail=f"缺少 python-pptx 依赖: {exc}") from exc

    prs = Presentation()
    prs.slide_width = 12192000
    prs.slide_height = 6858000
    blank_layout = prs.slide_layouts[6]
    for index, image_path in enumerate(image_paths):
        p = Path(image_path)
        if not p.exists() or not p.is_file():
            raise HTTPException(status_code=404, detail=f"图片文件不存在: {p}")
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(p), 0, 0, width=prs.slide_width, height=prs.slide_height)
        if notes and index < len(notes) and str(notes[index] or "").strip():
            slide.notes_slide.notes_text_frame.text = str(notes[index]).strip()

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return str(out)


def create_download_token(path: str) -> str:
    p = Path(path).resolve()
    allowed_root = default_create_ppt_output_dir().resolve()
    try:
        p.relative_to(allowed_root)
    except ValueError as exc:
        raise HTTPException(status_code=403, detail="PPT 文件不在允许下载目录内") from exc
    if not p.exists() or not p.is_file():
        raise HTTPException(status_code=404, detail="PPT 文件不存在")
    token = uuid.uuid4().hex
    token_dir = _lobster_root() / "_lobster_runtime" / "create_ppt"
    token_dir.mkdir(parents=True, exist_ok=True)
    token_file = token_dir / "download_tokens.json"
    try:
        raw = json.loads(token_file.read_text(encoding="utf-8")) if token_file.exists() else {}
    except Exception:
        raw = {}
    raw[token] = str(p)
    token_file.write_text(json.dumps(raw, ensure_ascii=False, indent=2), encoding="utf-8")
    return token


def resolve_download_token(token: str) -> Path:
    clean = (token or "").strip()
    if not clean:
        raise HTTPException(status_code=404, detail="下载 token 无效")
    token_file = _lobster_root() / "_lobster_runtime" / "create_ppt" / "download_tokens.json"
    try:
        raw = json.loads(token_file.read_text(encoding="utf-8")) if token_file.exists() else {}
    except Exception:
        raw = {}
    path = raw.get(clean)
    if not path:
        raise HTTPException(status_code=404, detail="下载 token 不存在或已失效")
    p = Path(path).resolve()
    allowed_root = default_create_ppt_output_dir().resolve()
    try:
        p.relative_to(allowed_root)
    except ValueError as exc:
        raise HTTPException(status_code=403, detail="PPT 文件不在允许下载目录内") from exc
    if not p.exists() or not p.is_file():
        raise HTTPException(status_code=404, detail="PPT 文件不存在")
    return p
